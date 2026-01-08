#!/usr/bin/env python3
"""
User Feedback Analyzer & Automated PowerPoint Generator
ØªØ·Ø¨ÙŠÙ‚ Flask Ù„ØªØ­Ù„ÙŠÙ„ Ù…Ù„Ø§Ø­Ø¸Ø§Øª Ø§Ù„Ø²Ø§Ø¦Ø± Ø§Ù„Ø³Ø±ÙŠ ÙˆØ¥Ù†Ø´Ø§Ø¡ ØªÙ‚Ø§Ø±ÙŠØ± PowerPoint
"""

from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import tempfile
import socket

app = Flask(__name__, static_folder='.')
CORS(app)

TEMPLATE_PATH = 'template.pptx'

def find_free_port():
    """Ø§Ù„Ø¨Ø­Ø« Ø¹Ù† Ø¨ÙˆØ±Øª ÙØ§Ø¶ÙŠ ØªÙ„Ù‚Ø§Ø¦ÙŠØ§Ù‹"""
    ports_to_try = [8080, 5000, 5001, 3000, 8000, 9000, 4000, 7000, 6000]
    for port in ports_to_try:
        try:
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
                s.bind(('', port))
                return port
        except OSError:
            continue
    # Ø¥Ø°Ø§ ÙƒÙ„ Ø§Ù„Ø¨ÙˆØ±ØªØ§Øª Ù…Ø´ØºÙˆÙ„Ø©ØŒ Ø§Ø³ØªØ®Ø¯Ù… Ø¨ÙˆØ±Øª Ø¹Ø´ÙˆØ§Ø¦ÙŠ
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(('', 0))
        return s.getsockname()[1]

@app.route('/')
def index():
    """Ø§Ù„ØµÙØ­Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©"""
    return send_from_directory('.', 'index.html')

@app.route('/generate-pptx', methods=['POST'])
def generate_pptx():
    """Ø¥Ù†Ø´Ø§Ø¡ Ø¨ÙˆØ±Ø¨ÙˆÙŠÙ†Øª Ù…Ù† Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª"""
    try:
        # Ø§Ø³ØªÙ„Ø§Ù… Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª
        data = request.get_json()
        
        if not data or 'categories' not in data or not data['categories']:
            return jsonify({'error': 'Ù„Ø§ ØªÙˆØ¬Ø¯ Ù…Ù„Ø§Ø­Ø¸Ø§Øª Ù„Ù„ØªØ­Ù„ÙŠÙ„'}), 400
        
        # ØªØ­ÙˆÙŠÙ„ categories Ø¥Ù„Ù‰ list Ø¥Ø°Ø§ ÙƒØ§Ù† dict
        categories_data = data.get('categories', [])
        
        if isinstance(categories_data, dict):
            # Ø¥Ø°Ø§ ÙƒØ§Ù† dictØŒ Ù†Ø­ÙˆÙ„Ù‡ Ù„Ù€ list Ù…Ù† Ø§Ù„Ù€ values
            # Ù„ÙƒÙ† ÙƒÙ„ value Ù‡Ùˆ category ÙƒØ§Ù…Ù„ (dict)
            categories = []
            for key, value in categories_data.items():
                if isinstance(value, dict):
                    # Ø¥Ø°Ø§ Ø§Ù„Ù€ value Ù‡Ùˆ category ÙƒØ§Ù…Ù„
                    categories.append(value)
                elif isinstance(value, list):
                    # Ø¥Ø°Ø§ Ø§Ù„Ù€ value Ù‡Ùˆ list of notes
                    categories.append({
                        'name': key,
                        'notes': value
                    })
        elif isinstance(categories_data, list):
            categories = categories_data
        else:
            categories = []
        
        # ÙØªØ­ Ø§Ù„Ù‚Ø§Ù„Ø¨
        if not os.path.exists(TEMPLATE_PATH):
            return jsonify({'error': f'Ø§Ù„Ù‚Ø§Ù„Ø¨ ØºÙŠØ± Ù…ÙˆØ¬ÙˆØ¯: {TEMPLATE_PATH}'}), 500
        
        prs = Presentation(TEMPLATE_PATH)
        
        # ØªÙ†Ø¸ÙŠÙ Ø§Ù„Ø´Ø±Ø§Ø¦Ø­ (Ø­Ø°Ù Ø§Ù„Ù…Ø­ØªÙˆÙ‰ØŒ Ø§Ù„Ø®Ù„ÙÙŠØ© ØªØ¨Ù‚Ù‰)
        for slide in prs.slides:
            shapes_to_delete = [shape for shape in slide.shapes]
            for shape in shapes_to_delete:
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                except:
                    pass
        
        # ØªØ¬Ù…ÙŠØ¹ Ø§Ù„Ù…Ù„Ø§Ø­Ø¸Ø§Øª Ø­Ø³Ø¨ Ø§Ù„ÙØ¦Ø© Ø§Ù„ÙØ±Ø¹ÙŠØ©
        by_subcategory = {}
        
        for category in categories:
            # Ø§Ù„ØªØ£ÙƒØ¯ Ù…Ù† Ø£Ù† category Ù‡Ùˆ dict
            if isinstance(category, str):
                continue
            
            category_name = category.get('name', 'ØºÙŠØ± Ù…Ø­Ø¯Ø¯')
            notes = category.get('notes', [])
            
            for note in notes:
                # Ø§Ù„ØªØ£ÙƒØ¯ Ù…Ù† Ø£Ù† note Ù‡Ùˆ dict
                if isinstance(note, str):
                    continue
                    
                subcategory = note.get('subCategory', 'ØºÙŠØ± Ù…Ø­Ø¯Ø¯')
                key = f"{category_name} ( {subcategory} )"
                
                if key not in by_subcategory:
                    by_subcategory[key] = []
                    
                by_subcategory[key].append({
                    'category': category_name,
                    'subCategory': subcategory,
                    'observation': note.get('observation', ''),
                    'repeatCount': note.get('repeatCount', 1)
                })
        
        if not by_subcategory:
            return jsonify({'error': 'Ù„Ù… ÙŠØªÙ… Ø§Ù„Ø¹Ø«ÙˆØ± Ø¹Ù„Ù‰ Ù…Ù„Ø§Ø­Ø¸Ø§Øª ØµØ§Ù„Ø­Ø©'}), 400
        
        slide_index = 0
        
        # === Ø´Ø±ÙŠØ­Ø© Ø§Ù„Ø¹Ù†ÙˆØ§Ù† ===
        slide = prs.slides[slide_index]
        slide_index += 1
        
        # Ø§Ù„Ø¹Ù†ÙˆØ§Ù† Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠ
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.0), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "Ø§Ù„Ø®Ø·Ø© Ø§Ù„ØªØµØ­ÙŠØ­ÙŠØ© Ù„Ù…Ù„Ø§Ø­Ø¸Ø§Øª Ø§Ù„Ø²Ø§Ø¦Ø± Ø§Ù„Ø³Ø±ÙŠ"
        
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.name = 'Calibri'
                run.font.bold = True
        
        # Ø§Ù„Ø¹Ù†ÙˆØ§Ù† Ø§Ù„ÙØ±Ø¹ÙŠ (Ø§Ù„Ù…Ù†Ø´Ø£Ø© ÙˆØ±Ù‚Ù… Ø§Ù„ØªØ°ÙƒØ±Ø©)
        subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.0), Inches(1.0))
        subtitle_frame = subtitle_box.text_frame
        subtitle_text = f"{data.get('hospital', 'ØºÙŠØ± Ù…Ø­Ø¯Ø¯')}\n{data.get('ticketId', 'ØºÙŠØ± Ù…Ø­Ø¯Ø¯')}"
        subtitle_frame.text = subtitle_text
        
        for paragraph in subtitle_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(24)
                run.font.name = 'Calibri'
        
        # === Ø´Ø±Ø§Ø¦Ø­ Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª ===
        for idx, (subcategory_full, notes) in enumerate(sorted(by_subcategory.items()), 1):
            
            # Ø§Ø³ØªØ®Ø¯Ø§Ù… Ø´Ø±ÙŠØ­Ø© Ù…Ù† Ø§Ù„Ù‚Ø§Ù„Ø¨
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                slide_index += 1
            else:
                # ØªÙƒØ±Ø§Ø± Ø§Ù„Ø´Ø±ÙŠØ­Ø© 2 Ø¥Ø°Ø§ Ø§Ù†ØªÙ‡Øª Ø§Ù„Ø´Ø±Ø§Ø¦Ø­
                template_slide = prs.slides[1]
                slide_layout = template_slide.slide_layout
                slide = prs.slides.add_slide(slide_layout)
                
                # ØªÙ†Ø¸ÙŠÙ Ø§Ù„Ø´Ø±ÙŠØ­Ø© Ø§Ù„Ø¬Ø¯ÙŠØ¯Ø©
                shapes_to_delete = [shape for shape in slide.shapes]
                for shape in shapes_to_delete:
                    try:
                        sp = shape.element
                        sp.getparent().remove(sp)
                    except:
                        pass
            
            # Ø­Ø³Ø§Ø¨ Ø§Ø±ØªÙØ§Ø¹ Ø§Ù„Ø¬Ø¯ÙˆÙ„ Ø­Ø³Ø¨ Ø¹Ø¯Ø¯ Ø§Ù„Ù…Ù„Ø§Ø­Ø¸Ø§Øª
            if len(notes) <= 3:
                height = Inches(1.5)
            elif len(notes) <= 6:
                height = Inches(2.5)
            elif len(notes) <= 9:
                height = Inches(3.5)
            else:
                height = Inches(4.0)
            
            # Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ø¬Ø¯ÙˆÙ„
            rows = len(notes) + 1
            cols = 3
            table_shape = slide.shapes.add_table(
                rows, cols, 
                Inches(1.5), Inches(2.0), 
                Inches(9.0), height
            )
            table = table_shape.table
            
            # Ø¹Ø±Ø¶ Ø§Ù„Ø£Ø¹Ù…Ø¯Ø©
            table.columns[0].width = Inches(0.8)  # Ø§Ù„Ø­Ø§Ù„Ø©
            table.columns[1].width = Inches(2.0)  # Ø§Ù„Ø®Ø·Ø©
            table.columns[2].width = Inches(8.5)  # Ø§Ù„Ù…Ù„Ø§Ø­Ø¸Ø©
            
            # Ø±Ø£Ø³ Ø§Ù„Ø¬Ø¯ÙˆÙ„ (Ø£Ø²Ø±Ù‚ ØºØ§Ù…Ù‚ + Ù†Øµ Ø£Ø¨ÙŠØ¶)
            headers = ['Ø§Ù„Ø­Ø§Ù„Ø©', 'Ø§Ù„Ø®Ø·Ø© Ø§Ù„ØªØµØ­ÙŠØ­ÙŠØ©', 'Ø§Ù„Ù…Ù„Ø§Ø­Ø¸Ø©']
            for col_idx, header_text in enumerate(headers):
                cell = table.rows[0].cells[col_idx]
                cell.text = header_text
                
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = 1
                
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(18)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
            
            # Ø¨ÙŠØ§Ù†Ø§Øª Ø§Ù„Ø¬Ø¯ÙˆÙ„ (Ø£Ù„ÙˆØ§Ù† Ù…ØªÙ†Ø§ÙˆØ¨Ø©)
            for row_idx, note in enumerate(notes, start=1):
                # Ù„ÙˆÙ† Ø§Ù„ØµÙ (Ù…ØªÙ†Ø§ÙˆØ¨)
                row_color = RGBColor(217, 225, 242) if row_idx % 2 == 1 else RGBColor(255, 255, 255)
                
                # Ø§Ù„Ø¹Ù…ÙˆØ¯ 1: Ø§Ù„Ø­Ø§Ù„Ø© (ÙØ§Ø±Øº)
                cell1 = table.rows[row_idx].cells[0]
                cell1.text = ""
                cell1.fill.solid()
                cell1.fill.fore_color.rgb = row_color
                cell1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Ø§Ù„Ø¹Ù…ÙˆØ¯ 2: Ø§Ù„Ø®Ø·Ø© Ø§Ù„ØªØµØ­ÙŠØ­ÙŠØ© (ÙØ§Ø±Øº)
                cell2 = table.rows[row_idx].cells[1]
                cell2.text = ""
                cell2.fill.solid()
                cell2.fill.fore_color.rgb = row_color
                cell2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Ø§Ù„Ø¹Ù…ÙˆØ¯ 3: Ø§Ù„Ù…Ù„Ø§Ø­Ø¸Ø©
                category = note['category']
                subcategory = note['subCategory']
                observation = note['observation']
                
                # Ø§Ù„Ù†Øµ Ø¨Ø¯ÙˆÙ† Ø¹Ø¯Ø¯ Ø§Ù„ØªÙƒØ±Ø§Ø±
                full_text = f"ÙÙŠ {category} ( {subcategory} ) {observation}"
                
                cell = table.rows[row_idx].cells[2]
                cell.text = full_text
                
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_color
                
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                
                cell.text_frame.word_wrap = True
                cell.vertical_anchor = 1
                cell.margin_left = Inches(0.1)
                cell.margin_right = Inches(0.1)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
        
        # Ø­Ø°Ù Ø§Ù„Ø´Ø±Ø§Ø¦Ø­ Ø§Ù„ÙØ§Ø¶ÙŠØ©
        total_slides_used = slide_index
        total_slides = len(prs.slides)
        
        if total_slides_used < total_slides:
            slides_to_delete = total_slides - total_slides_used
            
            for _ in range(slides_to_delete):
                rId = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[-1]
        
        # Ø­ÙØ¸ Ø§Ù„Ù…Ù„Ù
        output = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(output.name)
        output.close()
        
        return send_file(
            output.name,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f'ØªØ­Ù„ÙŠÙ„_Ø§Ù„Ø²Ø§Ø¦Ø±_Ø§Ù„Ø³Ø±ÙŠ_{data.get("ticketId", "")}.pptx'
        )
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    port = find_free_port()
    print("="*60)
    print("ğŸš€ User Feedback Analyzer & Automated PowerPoint Generator")
    print("="*60)
    print(f"âœ… Ø§Ù„Ø³ÙŠØ±ÙØ± ÙŠØ¹Ù…Ù„ Ø¹Ù„Ù‰: http://localhost:{port}")
    print(f"ğŸ“‚ Ø§ÙØªØ­ Ø§Ù„Ù…ØªØµÙØ­ Ø¹Ù„Ù‰: http://localhost:{port}")
    print("="*60)
    app.run(host='0.0.0.0', port=port, debug=True)
